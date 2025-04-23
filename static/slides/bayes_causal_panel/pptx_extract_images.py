import os
from pptx import Presentation
import io
from PIL import Image
import argparse

def extract_images(pptx_path, output_dir):
    # Create output directory if it doesn't exist
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
    
    # Open the presentation
    prs = Presentation(pptx_path)
    
    # Track image count for naming
    image_count = 0
    
    # Go through each slide
    for slide in prs.slides:
        # Go through each shape in the slide
        for shape in slide.shapes:
            # Check if the shape has an image
            if hasattr(shape, "image"):
                # Get the image info
                image = shape.image
                # Get the content type (e.g., 'image/jpeg')
                image_type = image.content_type
                # Get the extension from the content type
                if "png" in image_type:
                    ext = ".png"
                elif "jpeg" in image_type or "jpg" in image_type:
                    ext = ".jpg"
                else:
                    ext = ".img"  # Default extension
                
                # Save the image
                image_filename = f"image_{image_count}{ext}"
                image_path = os.path.join(output_dir, image_filename)
                
                # Convert image blob to file
                with open(image_path, "wb") as f:
                    f.write(image.blob)
                
                print(f"Saved {image_path}")
                image_count += 1

if __name__ == "__main__":
    # Set up argument parser
    parser = argparse.ArgumentParser(description='Extract images from a PowerPoint presentation')
    parser.add_argument('pptx_file', help='Path to the PowerPoint file')
    parser.add_argument('output_directory', help='Directory to save extracted images')
    
    # Parse arguments
    args = parser.parse_args()
    
    # Call the extraction function with command line arguments
    extract_images(args.pptx_file, args.output_directory)
    print(f"Extracted images to {args.output_directory}")